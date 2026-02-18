import logging
import os
from typing import Dict, Any

# Library imports handled inside functions to avoid hard dependency if not installed, 
# but requirements.txt has them now.
try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger("rog.loaders.office")

def load_docx(file_path: str) -> Dict[str, Any]:
    if not Document:
        logger.error("python-docx not installed.")
        return {"text": "", "error": "missing_dependency"}
        
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                full_text.append(" | ".join(row_data))
                
        return {
            "status": "success",
            "text": "\n".join(full_text),
            "metadata": {"type": "docx"}
        }
    except Exception as e:
        logger.error(f"Error reading DOCX {file_path}: {e}")
        return {"status": "error", "error": str(e), "text": ""}

def load_xlsx(file_path: str) -> Dict[str, Any]:
    if not openpyxl:
        logger.error("openpyxl not installed.")
        return {"text": "", "error": "missing_dependency"}
        
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        full_text = []
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            full_text.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                # Filter None values and convert to string
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    full_text.append(" | ".join(row_text))
                    
        return {
            "status": "success",
            "text": "\n".join(full_text),
            "metadata": {"type": "xlsx", "sheets": wb.sheetnames}
        }
    except Exception as e:
        logger.error(f"Error reading XLSX {file_path}: {e}")
        return {"status": "error", "error": str(e), "text": ""}

def load_pptx(file_path: str) -> Dict[str, Any]:
    if not Presentation:
        logger.error("python-pptx not installed.")
        return {"text": "", "error": "missing_dependency"}
        
    try:
        prs = Presentation(file_path)
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                full_text.append(f"--- Slide {i+1} ---")
                full_text.extend(slide_text)
                
        return {
            "status": "success",
            "text": "\n".join(full_text),
            "metadata": {"type": "pptx", "slides": len(prs.slides)}
        }
    except Exception as e:
        logger.error(f"Error reading PPTX {file_path}: {e}")
        return {"status": "error", "error": str(e), "text": ""}
